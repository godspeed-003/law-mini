import os
from typing import List, Dict, Any
import google.generativeai as gen_ai  
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings  # Updated import
from langchain_community.vectorstores import FAISS
import PyPDF2  # We'll use only this for PDF handling
from docx import Document
import pandas as pd
import sys
import tiktoken
import json
import xml.etree.ElementTree as ET
import yaml
import csv
import xlrd
import openpyxl
import pptx
import html
import re
import io
from PIL import Image
from google.cloud import vision
try:
    from pymupdf import fitz  # Try PyMuPDF first
    PYMUPDF_AVAILABLE = True
except ImportError:
    print("Warning: PyMuPDF not available, using PyPDF2 for PDF processing")
    PYMUPDF_AVAILABLE = False

import hashlib
import pickle

# Fix torch/transformers imports
try:
    import torch
    from transformers import AutoTokenizer, AutoModelForSequenceClassification, AutoModel
    from huggingface_hub import snapshot_download  # Add for offline caching
    TORCH_AVAILABLE = True
except ImportError as e:
    print(f"Warning: PyTorch/Transformers not available: {e}")
    TORCH_AVAILABLE = False

class DocumentProcessor:
    def __init__(self, api_key: str):
        self.api_key = api_key
        gen_ai.configure(api_key=api_key)  # Configure Gemini
        self.model = gen_ai.GenerativeModel('models/gemini-2.0-flash')
        # Use the multilingual-e5-small embedding model for better multilingual/document retrieval
        self.embeddings = HuggingFaceEmbeddings(model_name="intfloat/multilingual-e5-small")
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000,
            chunk_overlap=400,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
        # InLegalBERT setup with robust error handling and device management
        self.inlegalbert_tokenizer = None
        self.inlegalbert_statute_model = None
        self.inlegalbert_segmentation_model = None
        self.inlegalbert_judgment_model = None
        self.inlegalbert_label_maps = {
            "statute": None,
            "segmentation": None,
            "judgment": ["REJECTED", "ACCEPTED"]
        }
        self.inlegalbert_cache_dir = os.path.join(os.path.dirname(__file__), "inlegalbert_cache")
        os.makedirs(self.inlegalbert_cache_dir, exist_ok=True)
        # Initialize InLegalBERT models only if torch is available
        if TORCH_AVAILABLE:
            self._initialize_inlegalbert_models()
        else:
            print("Warning: PyTorch not available. InLegalBERT models will be disabled.")
            self._disable_inlegalbert()
        
        # Only keep OCR setup for images
        self.vision_client = vision.ImageAnnotatorClient()
        
    def perform_ocr(self, file_path: str) -> str:
        """Perform OCR on image files only."""
        try:
            return self._ocr_image(file_path)
        except Exception as e:
            raise ValueError(f"OCR Error: {str(e)}")

    def _ocr_image(self, image_path: str) -> str:
        """Process single image using Google Cloud Vision."""
        try:
            with open(image_path, "rb") as image_file:
                content = image_file.read()
            
            image = vision.Image(content=content)
            response = self.vision_client.text_detection(image=image)
            
            if response.error.message:
                raise Exception(response.error.message)
            
            texts = response.text_annotations
            return texts[0].description if texts else "No text detected"
            
        except Exception as e:
            raise ValueError(f"Image OCR error: {str(e)}")

    def _ocr_pil_image(self, pil_image: Image.Image) -> str:
        """Enhanced OCR for PIL Image with better error handling."""
        try:
            max_dimension = 3000
            if max(pil_image.size) > max_dimension:
                pil_image.thumbnail((max_dimension, max_dimension), Image.Resampling.LANCZOS)
            if pil_image.mode != 'RGB':
                pil_image = pil_image.convert('RGB')
            img_byte_arr = io.BytesIO()
            pil_image.save(img_byte_arr, format='PNG')
            img_byte_arr = img_byte_arr.getvalue()
            image = vision.Image(content=img_byte_arr)
            response = self.vision_client.text_detection(image=image)
            if response.error.message:
                raise Exception(response.error.message)
            texts = response.text_annotations
            raw_text = texts[0].description if texts else ""
            return self.clean_ocr_text(raw_text)
        except Exception as e:
            print(f"OCR failed for image: {e}")
            return ""

    def clean_ocr_text(self, text: str) -> str:
        """Clean OCR output text."""
        if not text:
            return ""
        
        # Basic cleaning
        text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with single space
        text = text.replace('\n', ' ').strip()  # Remove newlines
        text = re.sub(r'[^\w\s.,!?-]', '', text)  # Remove special characters except basic punctuation
        
        # Fix common OCR errors
        text = re.sub(r'(?<=[.,!?])\s*(?=[A-Z])', ' ', text)  # Ensure space after punctuation
        text = re.sub(r'\b0\b', 'O', text)  # Replace standalone '0' with 'O'
        text = re.sub(r'\bl\b', 'I', text)  # Replace standalone 'l' with 'I'
        
        return text.strip()

    def read_file(self, file_path: str) -> str:
        """Enhanced read_file that preserves natural document flow for PDFs."""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            if file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
                return self.perform_ocr(file_path)
            elif file_extension == '.pdf':
                return self._read_pdf_natural_flow(file_path)
            elif file_extension == '.docx':
                return self._read_docx(file_path)
            elif file_extension in ['.txt', '.text']:
                return self._read_txt(file_path)
            elif file_extension == '.csv':
                return self._read_csv(file_path)
            elif file_extension == '.json':
                return self._read_json(file_path)
            elif file_extension in ['.xml']:
                return self._read_xml(file_path)
            elif file_extension in ['.yml', '.yaml']:
                return self._read_yaml(file_path)
            elif file_extension == '.xls':
                return self._read_xls(file_path)
            elif file_extension == '.xlsx':
                return self._read_xlsx(file_path)
            elif file_extension == '.pptx':
                return self._read_pptx(file_path)
            elif file_extension in ['.html', '.htm']:
                return self._read_html(file_path)
            else:
                try:
                    return self._read_txt(file_path)
                except:
                    raise ValueError(f"Unsupported file type: {file_extension}")
        except Exception as e:
            raise ValueError(f"Error reading file {file_path}: {str(e)}")

    def _read_pdf_natural_flow(self, file_path: str) -> str:
        """Enhanced PDF reader using only PyPDF2."""
        try:
            return self._read_pdf_with_pypdf2(file_path)
        except Exception as e:
            print(f"PDF reading error: {e}")
            return "Error reading PDF file"

    def _read_pdf_with_pypdf2(self, file_path: str) -> str:
        """Read PDF with PyPDF2."""
        try:
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                text_parts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(text.strip())
                return "\n\n".join(text_parts) if text_parts else "No text content found in PDF"
        except Exception as e:
            print(f"PyPDF2 reading failed: {str(e)}")
            return f"Error reading PDF: {str(e)}"

    def _read_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _read_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def _read_csv(self, file_path: str) -> str:
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def _read_json(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            return json.dumps(data, indent=2)
    
    def _read_xml(self, file_path: str) -> str:
        tree = ET.parse(file_path)
        root = tree.getroot()
        return ET.tostring(root, encoding='unicode', method='xml')
    
    def _read_yaml(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = yaml.safe_load(file)
            return yaml.dump(data, default_flow_style=False)
    
    def _read_xls(self, file_path: str) -> str:
        workbook = xlrd.open_workbook(file_path)
        text = ""
        for sheet in workbook.sheets():
            text += f"Sheet: {sheet.name}\n"
            for row in range(sheet.nrows):
                text += "\t".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols)) + "\n"
        return text
    
    def _read_xlsx(self, file_path: str) -> str:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in workbook.sheetnames:
            text += f"Sheet: {sheet}\n"
            ws = workbook[sheet]
            for row in ws.rows:
                text += "\t".join(str(cell.value) for cell in row) + "\n"
        return text
    
    def _read_pptx(self, file_path: str) -> str:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            text += f"Slide {slide.slide_number}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _read_html(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
            # Remove HTML tags
            text = re.sub('<[^<]+?>', '', html_content)
            # Decode HTML entities
            text = html.unescape(text)
            return text
    
    def _get_type_from_extension(self, file_path: str) -> str:
        """Fallback method to detect file type from extension."""
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.txt': 'text/plain',
            '.csv': 'text/csv',
            '.json': 'application/json',
            '.xml': 'text/xml',
            '.yml': 'text/yaml',
            '.yaml': 'text/yaml',
            '.xls': 'application/vnd.ms-excel',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.html': 'text/html'
        }
        return type_map.get(ext, 'text/plain')
    
    def process_documents(self, file_paths: List[str]) -> FAISS:
        """Process multiple documents and create a FAISS vector store."""
        all_texts = []
        for file_path in file_paths:
            try:
                print(f"Processing file: {file_path}")  # Debug log
                text = self.read_file(file_path)
                # Remove minimum length requirement, just check if there's any content
                if text and text.strip():
                    print(f"Found {len(text.split())} words in {file_path}")  # Debug log
                    all_texts.append(text)
                else:
                    print(f"Warning: No content found in {file_path}")
            except Exception as e:
                print(f"Error processing {file_path}: {str(e)}")

        if not all_texts:
            print("No valid content found in any files. File contents:")  # Debug log
            for file_path in file_paths:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        print(f"\n{file_path} first 100 chars: {f.read(100)}")
                except:
                    print(f"Could not read {file_path}")
            raise ValueError("No valid text content found in any of the documents")

        # Split texts into chunks with minimum size decreased
        chunks = []
        for text in all_texts:
            text_chunks = self.text_splitter.split_text(text)
            # Reduce minimum chunk size to 1 word
            valid_chunks = [chunk for chunk in text_chunks if len(chunk.split()) >= 1]
            chunks.extend(valid_chunks)
            print(f"Created {len(valid_chunks)} chunks from text of length {len(text)}")  # Debug log

        if not chunks:
            raise ValueError("No valid chunks created from the documents")

        try:
            print(f"Creating embeddings for {len(chunks)} chunks...")  # Debug log
            # Create vector store with explicit error handling
            embeddings = self.embeddings.embed_documents(chunks[:1])  # Test with first chunk
            if not embeddings or len(embeddings[0]) == 0:
                raise ValueError("Failed to generate embeddings")
                
            vector_store = FAISS.from_texts(
                texts=chunks,
                embedding=self.embeddings,
                metadatas=[{"chunk": i} for i in range(len(chunks))]
            )
            print("Successfully created vector store")  # Debug log
            return vector_store
        except Exception as e:
            raise ValueError(f"Error creating vector store: {str(e)}")

    def generate_alternate_phrasings(self, user_question: str) -> list:
        """Use Gemini to generate alternate phrasings/meanings for the user query."""
        prompt = (
            "Given the following user question, generate a list of up to 5 alternate phrasings or possible meanings. "
            "These should include synonyms, rewordings, or clarifications that might help retrieve relevant information from a document database. "
            "Return the list as plain text, one alternate phrasing per line. Do not answer the question.\n\n"
            f"User Question: {user_question}\n"
        )
        chat_session = self.model.start_chat(history=[])
        generation_config = {
            "temperature": 0.3,
            "top_p": 0.8,
            "top_k": 40
        }
        try:
            response = chat_session.send_message(prompt, generation_config=generation_config)
            # Parse the response into a list of phrasings
            lines = [line.strip("-•* \n") for line in response.text.strip().splitlines() if line.strip()]
            # Filter out lines that are too short or identical to the original
            phrasings = [p for p in lines if p and p.lower() != user_question.lower()]
            # Always include the original question as the first phrasing
            return [user_question] + phrasings[:4]
        except Exception as e:
            return [user_question]

    def get_best_context_for_refinement(self, vector_store, user_question: str) -> list:
        """Aggregate top relevant chunks from similarity searches using alternate phrasings."""
        alternate_phrasings = self.generate_alternate_phrasings(user_question)
        seen_chunks = set()
        context_chunks = []
        # For each phrasing, perform similarity search and collect top chunks
        for phrasing in alternate_phrasings:
            docs = vector_store.similarity_search(phrasing, k=2)
            for doc in docs:
                # Use chunk index and filename as unique identifier
                chunk_id = (doc.metadata.get("source", "Unknown File"), doc.metadata.get("chunk", -1))
                if chunk_id not in seen_chunks:
                    seen_chunks.add(chunk_id)
                    context_chunks.append(doc)
                if len(context_chunks) >= 5:
                    break
            if len(context_chunks) >= 5:
                break
        return context_chunks

    def merge_relevant_chunks(self, docs, query):
        """
        Merge consecutive relevant chunks from the same document if they mention the same course/unit.
        """
        merged = []
        last_filename = None
        buffer = []
        keywords = self.extract_keywords_from_query(query)
        for doc in docs:
            filename = doc.metadata.get("source", "Unknown File")
            chunk_text = doc.page_content.strip()
            # Check if chunk is relevant to the query (contains any keyword)
            if any(kw.lower() in chunk_text.lower() for kw in keywords):
                if filename == last_filename:
                    buffer.append(chunk_text)
                else:
                    if buffer:
                        merged.append({"filename": last_filename, "content": "\n".join(buffer)})
                    buffer = [chunk_text]
                    last_filename = filename
            else:
                if buffer:
                    merged.append({"filename": last_filename, "content": "\n".join(buffer)})
                    buffer = []
                    last_filename = None
        if buffer:
            merged.append({"filename": last_filename, "content": "\n".join(buffer)})
        return merged if merged else [{"filename": doc.metadata.get("source", "Unknown File"), "content": doc.page_content.strip()} for doc in docs]

    def extract_keywords_from_query(self, query):
        """
        Extract course/unit/subject keywords from the query for better chunk merging.
        """
        # Simple heuristic: split on 'and', ',', or 'of', remove stopwords
        words = re.split(r'\band\b|,|of|for|in|the|to|with', query, flags=re.IGNORECASE)
        keywords = [w.strip() for w in words if len(w.strip()) > 2]
        return keywords if keywords else [query]

    def refine_prompt(self, user_question: str, relevant_docs: list) -> str:
        """Refine the user question using Gemini and relevant document context."""
        context_chunks = []
        for doc in relevant_docs:
            filename = doc.metadata.get("source", "Unknown File")
            chunk_text = doc.page_content
            context_chunks.append(f"Filename: {os.path.basename(filename)}\nChunk: {chunk_text}")

        context_str = "\n\n".join(context_chunks)

        refinement_prompt = (
            "You are an expert assistant tasked with improving user questions for a document Q&A system. "
            "Your job is to rewrite the user's question to be more precise and targeted to the available context.\n\n"
            "Requirements:\n"
            "1. Keep the core intent of the question unchanged\n"
            "2. Use proper terminology found in the context documents\n"
            "3. Make the question specific and unambiguous\n"
            "4. Fix any spelling or grammar issues\n"
            "5. If asking about multiple topics, clearly separate them\n"
            "6. Focus on information that appears in the context\n"
            "7. Return ONLY the refined question - no explanations or other text\n"
            "8. It is not your job if the user question doesn't match any of the processed documents, in that case just refine the question no matter how irrelevant and then do not return anything saying 'no relevant information found'. Your work is to refine it no matter the relevance, so just do it, you may show a disclaimer in final answer that no such information was found and that this is solely based on gemini's knowledge so it may be inaccurate.\n"
            "9. Keep the question natural and conversational\n"
            "10. Do not include phrases like 'Based on the context' or 'In the documents'\n"
            "11. Never say the question cannot be answered - just improve its structure\n\n"
            f"Original Question: {user_question}\n\n"
            f"Available Context:\n{context_str if context_str else '[No context available]'}\n\n"
            "Refined Question:"
        )

        chat_session = self.model.start_chat(history=[])
        generation_config = {
            "temperature": 0.2,
            "top_p": 0.8,
            "top_k": 40
        }
        try:
            response = chat_session.send_message(refinement_prompt, generation_config=generation_config)
            refined = response.text.strip()
            # Always return a refinement prompt, never just a disclaimer
            if not refined or "cannot refine" in refined.lower() or "not relevant" in refined.lower():
                # Try to at least clarify or rephrase the original question
                correction_prompt = (
                    "Rewrite the following question to be as clear and specific as possible, correcting spelling and grammar. "
                    "If it is ambiguous, make a best effort to clarify the intent. Return only the improved question.\n\n"
                    f"Question: {user_question}\n"
                )
                response2 = chat_session.send_message(correction_prompt, generation_config=generation_config)
                refined2 = response2.text.strip()
                return refined2 if refined2 else user_question
            return refined
        except Exception as e:
            return user_question

    def direct_gemini_query(self, question: str) -> str:
        """
        Query Gemini directly without RAG context for general knowledge questions.
        Always answer the user's question, even if unrelated to the documents.
        If the answer is unrelated, prepend a disclaimer.
        """
        prompt = (
            "You are a helpful AI assistant. Always answer the user's question, even if it is unrelated to the provided documents. "
            "If the answer is unrelated to the documents, start your answer with this disclaimer: "
            "'Disclaimer: The following answer is based solely on Gemini's general knowledge and not on the provided documents.'\n\n"
            "Answer the user's question as directly and thoroughly as possible, in a detailed and natural style. "
            "Do NOT provide any extra context or mention the documents unless specifically asked.\n\n"
            "If the question is inappropriate or potentially harmful, respond with 'I cannot assist with that.'\n\n"
            f"User Question: {question}\n\n"
            "Your answer:"
        )
        chat_session = self.model.start_chat(history=[])
        try:
            response = chat_session.send_message(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Error from Gemini: {str(e)}"

    def get_inlegalbert_analysis(self, query):
        """Run all three InLegalBERT tasks and return formatted analysis (with disk caching)."""
        # Check if models are loaded
        if not self.inlegalbert_tokenizer or not self.inlegalbert_statute_model:
            return "Legal Expert Analysis: InLegalBERT models not available. Please ensure the models are properly installed."
        cache_key = f"inlegalbert::{query}"
        cached = self._cache_inlegalbert(cache_key)
        if cached:
            return cached

        tokenizer = self.inlegalbert_tokenizer
        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        analysis = {}

        # Statute Identification (multi-label)
        try:
            if self.inlegalbert_statute_model is not None:
                inputs = tokenizer(query, return_tensors="pt", truncation=True, padding=True, max_length=512)
                model = self.inlegalbert_statute_model.to(device)
                with torch.no_grad():
                    outputs = model(**{k: v.to(device) for k, v in inputs.items()})
                logits = outputs.logits.cpu().numpy()[0]
                statutes = []
                if self.inlegalbert_label_maps["statute"] is not None:
                    for i, score in enumerate(logits):
                        if score > 0:
                            label = self.inlegalbert_label_maps["statute"].get(i, f"Statute_{i}")
                            statutes.append(label)
                else:
                    for i, score in enumerate(logits):
                        if score > 0:
                            statutes.append(f"Statute_Category_{i}")
                analysis["statutes"] = statutes if statutes else ["No relevant statutes identified"]
            else:
                analysis["statutes"] = ["Statute model not loaded"]
        except Exception as e:
            analysis["statutes"] = [f"Statute analysis error: {str(e)}"]

        # Segmentation (sentence tagging)
        try:
            if self.inlegalbert_segmentation_model is not None:
                model = self.inlegalbert_segmentation_model.to(device)
                sentences = re.split(r'(?<=[.!?])\s+', query)
                seg_labels = []
                for sent in sentences:
                    if not sent.strip():
                        continue
                    try:
                        inputs = tokenizer(sent, return_tensors="pt", truncation=True, padding=True, max_length=128)
                        with torch.no_grad():
                            outputs = model(**{k: v.to(device) for k, v in inputs.items()})
                        pred = outputs.logits.argmax(dim=-1).item()
                        if self.inlegalbert_label_maps["segmentation"] is not None and pred < len(self.inlegalbert_label_maps["segmentation"]):
                            label = self.inlegalbert_label_maps["segmentation"][pred]
                        else:
                            segment_labels = ["PREAMBLE", "FACTS", "ISSUES", "ARGUMENTS", "STATUTE", "PRECEDENT", "RATIO"]
                            label = segment_labels[pred] if pred < len(segment_labels) else f"Segment_{pred}"
                        seg_labels.append((sent.strip(), label))
                    except Exception as sent_error:
                        seg_labels.append((sent.strip(), f"Error: {str(sent_error)}"))
                analysis["segmentation"] = seg_labels if seg_labels else [("No sentences to analyze", "NONE")]
            else:
                analysis["segmentation"] = [("Segmentation model not loaded", "ERROR")]
        except Exception as e:
            analysis["segmentation"] = [("Segmentation analysis error", str(e))]

        # Judgment Prediction (binary)
        try:
            if self.inlegalbert_judgment_model is not None:
                model = self.inlegalbert_judgment_model.to(device)
                inputs = tokenizer(query, return_tensors="pt", truncation=True, padding=True, max_length=512)
                with torch.no_grad():
                    outputs = model(**{k: v.to(device) for k, v in inputs.items()})
                pred = outputs.logits.argmax(dim=-1).item()
                if pred < len(self.inlegalbert_label_maps["judgment"]):
                    label = self.inlegalbert_label_maps["judgment"][pred]
                else:
                    label = "UNKNOWN"
                analysis["judgment"] = label
            else:
                analysis["judgment"] = "Judgment model not loaded"
        except Exception as e:
            analysis["judgment"] = f"Judgment analysis error: {str(e)}"

        # Format for Gemini and UI
        formatted = "Legal Expert Analysis (InLegalBERT):\n"
        formatted += "\n1. Relevant Statutes Identified:\n"
        if analysis["statutes"]:
            formatted += "   - " + "\n   - ".join(analysis["statutes"])
        else:
            formatted += "   None found"
        formatted += "\n\n2. Document Segmentation:\n"
        for sent, label in analysis["segmentation"]:
            formatted += f"   - [{label}] {sent}\n"
        formatted += "\n3. Judgment Prediction:\n"
        formatted += f"   - The claim/petition is likely to be: {analysis['judgment']}\n"
        self._cache_inlegalbert(cache_key, formatted)
        return formatted

    def create_qa_chain(self, vector_store: FAISS):
        """
        Create a hybrid question-answering functionality using both RAG and direct Gemini queries,
        and return InLegalBERT context for UI.
        Gemini answers should be detailed and natural, not split into 'information from documents' and 'supplemental' sections.
        """
        try:
            def qa_function(user_question: str) -> Dict[str, Any]:
                inlegalbert_context = self.get_inlegalbert_analysis(user_question)
                context_docs = self.get_best_context_for_refinement(vector_store, user_question)
                refined_prompt = self.refine_prompt(user_question, context_docs)
                answer_docs = vector_store.similarity_search(refined_prompt, k=12)
                merged_chunks = self.merge_relevant_chunks(answer_docs, refined_prompt)
                merged_chunks = merged_chunks[:8]

                # If no relevant content, use direct Gemini query (with disclaimer if unrelated)
                if not merged_chunks:
                    answer = self.direct_gemini_query(refined_prompt)
                    return {
                        "result": answer,
                        "refined_prompt": refined_prompt,
                        "source_documents": [],
                        "expert": inlegalbert_context or ""
                    }

                # If we have relevant content, provide a detailed answer using Gemini, but do not show context blocks
                context = ""
                for chunk in merged_chunks:
                    context += f"\n{chunk['content']}\n"

                prompt = (
                    "You are an expert assistant for a Retrieval-Augmented Generation (RAG) system. "
                    "You are also provided with a Legal Expert Analysis from InLegalBERT. "
                    "Always use the Legal Expert Analysis for legal reasoning if available. "
                    "Answer the user's question in a detailed, natural, and conversational style. "
                    "Do NOT split your answer into 'information from documents' and 'supplemental' sections. "
                    "Do NOT show the context or mention the documents unless specifically asked. "
                    "If the answer is unrelated to the documents, start with this disclaimer: "
                    "'Disclaimer: The following answer is based solely on Gemini's general knowledge and not on the provided documents.'\n\n"
                    f"{inlegalbert_context}\n\n"
                    f"Context:\n{context}\n"
                    f"Question (refined): {refined_prompt}\n\n"
                    "Provide a long, detailed, and natural answer to the user's question."
                )

                chat_session = self.model.start_chat(history=[])
                try:
                    response = chat_session.send_message(prompt)
                    answer = response.text
                except Exception as e:
                    answer = f"Error from Gemini: {str(e)}"

                sources = []
                for chunk in merged_chunks:
                    sources.append(type("SourceDoc", (), {
                        "page_content": chunk["content"],
                        "filename": os.path.basename(chunk["filename"])
                    })())

                return {
                    "result": answer,
                    "refined_prompt": refined_prompt,
                    "source_documents": sources,
                    "expert": inlegalbert_context or ""
                }
            return qa_function
        except Exception as e:
            raise ValueError(f"Error creating QA function: {str(e)}")

    def _initialize_inlegalbert_models(self):
        """Initialize InLegalBERT models with proper error handling."""
        try:
            # Initialize tokenizer and models with cached downloads
            try:
                # Try downloading models to cache first
                model_paths = {
                    "base": snapshot_download("law-ai/InLegalBERT", local_dir=self.inlegalbert_cache_dir),
                    "ilsi": snapshot_download("law-ai/InLegalBERT-ilsi", local_dir=self.inlegalbert_cache_dir),
                    "iss": snapshot_download("law-ai/InLegalBERT-iss", local_dir=self.inlegalbert_cache_dir), 
                    "ildc": snapshot_download("law-ai/InLegalBERT-ildc", local_dir=self.inlegalbert_cache_dir)
                }
                
                # Now load from cache
                self.inlegalbert_tokenizer = AutoTokenizer.from_pretrained(
                    model_paths["base"],
                    local_files_only=True
                )
                self.inlegalbert_statute_model = AutoModelForSequenceClassification.from_pretrained(
                    model_paths["ilsi"],
                    local_files_only=True
                )
                self.inlegalbert_segmentation_model = AutoModelForSequenceClassification.from_pretrained(
                    model_paths["iss"],
                    num_labels=7,
                    local_files_only=True
                )
                self.inlegalbert_judgment_model = AutoModelForSequenceClassification.from_pretrained(
                    model_paths["ildc"],
                    num_labels=2,
                    local_files_only=True
                )
                print("InLegalBERT models loaded from cache")
                
            except Exception as cache_error:
                print(f"Cache loading failed: {cache_error}, trying direct download...")
                # Fallback to direct loading
                self.inlegalbert_tokenizer = AutoTokenizer.from_pretrained("law-ai/InLegalBERT")
                self.inlegalbert_statute_model = AutoModelForSequenceClassification.from_pretrained("law-ai/InLegalBERT-ilsi")
                self.inlegalbert_segmentation_model = AutoModelForSequenceClassification.from_pretrained("law-ai/InLegalBERT-iss", num_labels=7)
                self.inlegalbert_judgment_model = AutoModelForSequenceClassification.from_pretrained("law-ai/InLegalBERT-ildc", num_labels=2)
            
            self.inlegalbert_label_maps["statute"] = self._load_label_map("ilsi_labels.pkl")
            self.inlegalbert_label_maps["segmentation"] = self._load_label_map("iss_labels.pkl")
            print("InLegalBERT models loaded successfully")
            
        except Exception as e:
            print(f"Warning: InLegalBERT models could not be loaded: {e}")
            self._disable_inlegalbert()

    def _disable_inlegalbert(self):
        """Disable InLegalBERT functionality."""
        self.inlegalbert_tokenizer = None
        self.inlegalbert_statute_model = None
        self.inlegalbert_segmentation_model = None
        self.inlegalbert_judgment_model = None
        self.inlegalbert_label_maps["statute"] = None
        self.inlegalbert_label_maps["segmentation"] = None

    def _load_label_map(self, filename: str) -> Dict:
        """Load label map from pickle file or create default mapping."""
        try:
            filepath = os.path.join(self.inlegalbert_cache_dir, filename)
            if os.path.exists(filepath):
                with open(filepath, 'rb') as f:
                    return pickle.load(f)
            else:
                # Create default mappings
                if "ilsi" in filename:
                    return {i: f"Statute_Category_{i}" for i in range(50)}
                elif "iss" in filename:
                    return ["PREAMBLE", "FACTS", "ISSUES", "ARGUMENTS", "STATUTE", "PRECEDENT", "RATIO"]
                return {}
        except Exception as e:
            print(f"Warning: Could not load label map {filename}: {e}")
            return {}

    def _cache_inlegalbert(self, cache_key: str, data: str = None) -> str:
        """Simple disk cache for InLegalBERT results."""
        try:
            cache_file = os.path.join(
                self.inlegalbert_cache_dir, 
                f"{hashlib.md5(cache_key.encode()).hexdigest()}.txt"
            )
            if data is None:  # Read from cache
                if os.path.exists(cache_file):
                    with open(cache_file, 'r', encoding='utf-8') as f:
                        return f.read()
                return None
            else:  # Write to cache
                with open(cache_file, 'w', encoding='utf-8') as f:
                    f.write(data)
                return data
        except Exception as e:
            print(f"Warning: Cache operation failed: {e}")
            return None if data is None else data